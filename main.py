import base64
import importlib
from pathlib import Path
from pptx import Presentation
from astrbot.api.event import filter, AstrMessageEvent
from astrbot.api.star import Context, Star
from astrbot.api import logger, llm_tool, AstrBotConfig
from astrbot.api.provider import ProviderRequest
from astrbot.core.message.message_event_result import MessageChain, MessageEventResult
from astrbot.core.platform.message_type import MessageType
from astrbot.core.message.components import At, Reply
import astrbot.api.message_components as Comp

# 导入底层生成器
from .office_generator import OfficeGenerator
from .utils import format_file_size

TEXT_SUFFIXES = frozenset(
    {
        ".txt",
        ".md",
        ".py",
        ".js",
        ".ts",
        ".json",
        ".csv",
        ".html",
        ".css",
        ".yaml",
        ".yml",
        ".xml",
        ".sql",
        ".sh",
        ".bat",
        ".c",
        ".cpp",
        ".java",
    }
)

MAX_TEXT_READ = 200 * 1024  # 200 KB


class FileOperationPlugin(Star):
    """基于工具调用的智能文件管理插件"""

    def __init__(self, context: Context, config: AstrBotConfig):
        super().__init__(context)
        self.config = config

        self.plugin_data_path = Path(__file__).parent / "files"
        self.plugin_data_path.mkdir(parents=True, exist_ok=True)

        self.office_gen = OfficeGenerator(self.plugin_data_path)

        self.FILE_TOOLS = ["list_files", "read_file", "write_file", "delete_file"]
        self._office_libs = self._check_office_libs()
        logger.info(f"[文件管理] 插件加载完成。数据目录: {self.plugin_data_path}")

    def _check_permission(self, event: AstrMessageEvent) -> bool:
        """检查用户权限"""
        logger.info("正在检查用户权限")
        perm_cfg = self.config.get("permission_settings", {})

        # 管理员检查
        if perm_cfg.get("require_admin", False) and not event.is_admin():
            return False

        # 白名单检查
        whitelist = perm_cfg.get("whitelist_users", [])
        if whitelist:
            user_id = str(event.get_sender_id())
            if user_id not in [str(u) for u in whitelist]:
                return False

        return True

    def _is_bot_mentioned(self, event: AstrMessageEvent) -> bool:
        """检查是否被@/回复"""
        try:
            bot_id = str(event.message_obj.self_id)
            for segment in event.message_obj.message:
                if isinstance(segment, At) or isinstance(segment, Reply):
                    target_id = getattr(segment, "qq", None) or getattr(
                        segment, "target", None
                    )
                    if target_id and str(target_id) == bot_id:
                        return True
            return False
        except Exception as e:
            logger.error(f"未知错误{e}")
            return False

    def _validate_path(self, filename: str) -> tuple[bool, Path, str]:
        """
        验证文件路径安全性
        返回: (是否有效, 文件路径, 错误信息)
        """
        file_path = self.plugin_data_path / filename
        try:
            resolved = file_path.resolve()
            base = self.plugin_data_path.resolve()
            if not resolved.is_relative_to(base):
                return False, file_path, "非法路径：禁止访问工作区外的文件"
            return True, file_path, ""
        except Exception as e:
            return False, file_path, f"路径解析失败: {e}"

    def _check_office_libs(self) -> dict:
        """检查并缓存 Office 库的可用性"""
        libs = {}
        lib_names = {
            "docx": "python-docx",
            "openpyxl": "openpyxl",
            "pptx": "python-pptx",
        }
        for module_name, package_name in lib_names.items():
            try:
                libs[module_name] = importlib.import_module(module_name)
                logger.debug(f"[文件管理] {package_name} 已加载")
            except ImportError:
                libs[module_name] = None
                logger.warning(f"[文件管理] {package_name} 未安装")
        return libs

    def _get_max_file_size(self) -> int:
        """获取最大文件大小（字节）"""
        mb = self.config.get("file_settings", {}).get("max_file_size_mb", 50)
        return mb * 1024 * 1024

    def _get_max_read_text_size(self) -> int:
        """获取文本预览最大大小（字节）"""
        kb = self.config.get("file_settings", {}).get("max_read_text_kb", 100)
        return kb * 1024

    def _get_allowed_extensions(self) -> set:
        """获取允许的扩展名集合"""
        extensions = self.config.get("file_settings", {}).get("allowed_extensions", [])
        if not extensions:
            return set()  # 空集合表示允许所有
        # 确保扩展名以点开头
        return {f".{ext.lower().lstrip('.')}" for ext in extensions}

    def _is_extension_allowed(self, filename: str) -> bool:
        """检查文件扩展名是否允许"""
        allowed = self._get_allowed_extensions()
        if not allowed:  # 空集合 = 允许所有
            return True
        suffix = Path(filename).suffix.lower()
        return suffix in allowed

    @filter.on_llm_request()
    async def before_llm_chat(self, event: AstrMessageEvent, req: ProviderRequest):
        """动态控制工具可见性"""
        trigger_cfg = self.config.get("trigger_settings", {})

        is_group = event.message_obj.type == MessageType.GROUP_MESSAGE
        should_expose = True

        # 权限拦截
        if not self._check_permission(event):
            should_expose = False
        # 群聊@/回复拦截
        elif (
            is_group
            and trigger_cfg.get("require_at_in_group", True)
            and not self._is_bot_mentioned(event)
        ):
            should_expose = False

        if not should_expose and req.func_tool:
            for tool_name in self.FILE_TOOLS:
                req.func_tool.remove_tool(tool_name)

    @llm_tool(name="list_files")
    async def list_files(self, event: AstrMessageEvent):
        """列出机器人文件库中的所有文件。"""

        if not self._check_permission(event):
            await event.send(MessageChain().message("❌ 拒绝访问：权限不足"))
            return "拒绝访问：权限不足"
        try:
            office_suffixes = {".docx", ".xlsx", ".pptx"}
            files = [
                f
                for f in self.plugin_data_path.glob("*")
                if f.is_file() and f.suffix.lower() in office_suffixes
            ]
            if not files:
                msg = "文件库当前没有 Office 文件"
                await event.send(MessageChain().message(msg))
                return msg

            files.sort(key=lambda x: x.stat().st_mtime, reverse=True)
            res = ["📂 机器人工作区 Office 文件列表："]
            for f in files:
                res.append(f"- {f.name} ({format_file_size(f.stat().st_size)})")

            result = "\n".join(res)
            await event.send(MessageChain().message(result))
            return result
        except Exception as e:
            logger.error(f"获取列表失败: {e}")
            await event.send(MessageChain().message("获取列表失败"))
            return f"获取列表失败: {e}"

    @llm_tool(name="read_file")
    async def read_file(self, event: AstrMessageEvent, filename: str) -> str | None:
        """读取并查看文件内容。"""
        if not self._check_permission(event):
            await event.send(MessageChain().message("❌ 拒绝访问：权限不足"))
            return ""
        valid, file_path, error = self._validate_path(filename)
        if not valid:
            await event.send(MessageChain().message(f"❌ {error}"))
            return error
        if not file_path.exists():
            await event.send(MessageChain().message("文件不存在，请检查"))
            return ""
        if not self._is_extension_allowed(filename):
            await event.send(MessageChain().message("❌ 不支持的文件类型"))
            return "错误：不支持的文件类型"

        file_size = file_path.stat().st_size
        max_size = self._get_max_file_size()
        if file_size > max_size:
            size_str = format_file_size(file_size)
            max_str = format_file_size(max_size)
            await event.send(
                MessageChain().message(f"❌ 文件过大 ({size_str})，限制 {max_str}")
            )
            return f"错误：文件大小 {size_str} 超过限制 {max_str}"
        try:
            suffix = file_path.suffix.lower()
                # 文本文件：使用流式读取并限制最大读取量以防止内存耗尽
            if suffix in {".txt", ".md", ".json", ".csv", ".log", ".py", ".js", ".html", ".css", ".xml", ".yaml", ".yml"}:
                max_text = self._get_max_read_text_size()
                try:
                    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                        content = f.read(max_text + 1)

                    if len(content) > max_text:
                        content = content[:max_text]
                        truncated = True
                    else:
                        truncated = False

                    result = f"📄 文件: {filename}\n"
                    result += f"📏 大小: {format_file_size(file_size)}\n"
                    if truncated:
                        result += f"⚠️ 内容已截断（显示前 {format_file_size(max_text)}）\n"
                    result += f"{'─' * 30}\n{content}"

                    await event.send(MessageChain().message(result[:100]))
                    return result

                except Exception as e:
                    logger.error(f"读取文件失败: {e}")
                    return f"读取失败: {e}"

            # Office 文件：尝试提取文本（若未安装对应解析库，则提示为二进制）
            office_suffixes = {".docx", ".xlsx", ".pptx"}
            if suffix in office_suffixes:
                extracted = None
                try:  # ← 添加 try 块
                    if suffix == ".docx" and self._office_libs.get("docx"):
                        from docx import Document

                        doc = Document(file_path)
                        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                        extracted = "\n".join(paragraphs)

                    elif suffix == ".xlsx" and self._office_libs.get("openpyxl"):
                        from openpyxl import load_workbook

                        wb = load_workbook(file_path, read_only=True, data_only=True)
                        texts = []
                        for ws in wb.worksheets:
                            for row in ws.iter_rows(values_only=True):
                                texts.append(
                                    "\t".join("" if v is None else str(v) for v in row)
                                )
                                if len("\n".join(texts)) > MAX_TEXT_READ:
                                    break
                            if len("\n".join(texts)) > MAX_TEXT_READ:
                                break
                        extracted = "\n".join(texts)
                    elif suffix == ".pptx" and self._office_libs.get("pptx"):
                        prs = Presentation(file_path)
                        texts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    texts.append(shape.text)
                                if len("\n".join(texts)) > MAX_TEXT_READ:
                                    break
                            if len("\n".join(texts)) > MAX_TEXT_READ:
                                break
                        extracted = "\n".join(texts)
                except Exception as exc:
                    logger.warning(f"Office 文本提取失败: {exc}", exc_info=True)

                if extracted:
                    if len(extracted) > MAX_TEXT_READ:
                        extracted = extracted[:MAX_TEXT_READ] + "\n\n...（已截断）..."
                    await event.send(MessageChain().message(f"提取内容:\n{extracted}"))
                    return extracted

                await event.send(
                    MessageChain().message(
                        "该文件为二进制格式或未安装解析库，无法直接读取。"
                    )
                )
                return "该文件为二进制格式，无法直接读取。"
        except Exception as e:
            await event.send(MessageChain().message("文件不存在，请检查"))
            return f"读取失败: {e}"

    @llm_tool(name="write_file")
    async def write_file(
        self,
        event: AstrMessageEvent,
        filename: str,
        content: str,
        file_type: str = "text",
    ):
        """在机器人工作区中创建或更新文件（仅支持 Office 文件）。"""
        filename = Path(filename).name
        if not self._check_permission(event):
            await event.send(MessageChain().message("❌ 拒绝访问：权限不足"))

        file_type_lower = file_type.lower()
        # 目前仅支持 Office 文件的生成
        if file_type_lower not in ["word", "excel", "powerpoint"]:
            await event.send(
                MessageChain().message(
                    "❌ 错误：当前仅支持 Office 文件（word/excel/powerpoint）生成。"
                )
            )

        if not self.config.get("feature_settings", {}).get("enable_office_files", True):
            await event.send(
                MessageChain().message("错误：当前配置禁用了 Office 文件生成功能。")
            )

        def format_file_size(size_bytes: int) -> str:
            """格式化文件大小为可读格式"""
            if size_bytes < 1024:
                return f"{size_bytes} B"
            elif size_bytes < 1024 * 1024:
                return f"{size_bytes / 1024:.1f} KB"
            elif size_bytes < 1024 * 1024 * 1024:
                return f"{size_bytes / (1024 * 1024):.1f} MB"
            else:
                return f"{size_bytes / (1024 * 1024 * 1024):.2f} GB"

        file_info = {
            "type": file_type_lower,
            "filename": filename,
            "content": content,
        }
        try:
            file_path = await self.office_gen.generate(
                event, file_info["type"], file_info
            )
            if file_path and file_path.exists():
                file_size = file_path.stat().st_size
                max_size = self._get_max_file_size()

                if file_size > max_size:
                    # 删除过大的文件
                    file_path.unlink()
                    size_str = format_file_size(file_size)
                    max_str = format_file_size(max_size)
                    await event.send(
                        MessageChain().message(f"❌ 生成的文件过大 ({size_str})，超过限制 {max_str}")
                    )

                with open(file_path, "rb") as f:
                    b64_str = base64.b64encode(f.read()).decode("utf-8")

                use_reply = self.config.get("trigger_settings", {}).get(
                    "reply_to_user", True
                )
                is_at = Comp.At(qq=event.get_sender_id()) if use_reply else None
                chain = [
                    Comp.Plain(f"✅ 文件已处理成功：{file_path.name}"),
                    is_at,
                    Comp.File(file=f"base64://{b64_str}", name=file_path.name),
                ]

                yield event.chain_result(chain)
                await event.send(
                    MessageChain().message(f"✅ 文件已处理成功：{file_path.name}")
                )
        except Exception as e:
            await event.send(MessageChain().message(f"文件操作异常: {e}"))

    @llm_tool(name="delete_file")
    async def delete_file(self, event: AstrMessageEvent, filename: str) -> str:
        """从工作区中永久删除指定文件。"""

        if not self._check_permission(event):
            await event.send(MessageChain().message("❌ 拒绝访问：权限不足"))
            return ""
        valid, file_path, error = self._validate_path(filename)
        if not valid:
            return f"❌ {error}"

        if file_path.exists():
            try:
                file_path.unlink(missing_ok=True)
                await event.send(
                    MessageChain().message(f"成功：文件 '{filename}' 已删除。")
                )
                return ""
            except IsADirectoryError:
                await event.send(MessageChain().message(f"'{filename}'是目录,拒绝删除"))
                return ""
            except PermissionError:
                await event.send(MessageChain().message("❌ 权限不足，无法删除文件"))
                return ""
            except Exception as e:
                logger.error(f"删除文件时发生错误{e}")
                await event.send(MessageChain().message(f"删除文件时发生错误{e}"))
                return ""
        await event.send(MessageChain().message(f"错误：找不到文件 '{filename}'"))
        return ""

    @filter.command("fileinfo")
    async def fileinfo(self, event: AstrMessageEvent):
        """显示文件管理工具的运行信息"""
        yield event.plain_result(
            "📂 AstrBot 文件操作工具\n"
            f"工作目录: {self.plugin_data_path}\n"
            f"回复模式: {'开启' if self.config.get('trigger_settings', {}).get('reply_to_user') else '关闭'}"
        )
