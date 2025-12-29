"""
受限代码执行模块

安全地执行 LLM 生成的 Python 代码来创建 Office 文件。
只允许导入特定的安全库，文件操作限制在工作目录内。
"""

import asyncio
import re
import traceback
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

# 允许导入的模块白名单
ALLOWED_IMPORTS = {
    # Office 库
    "openpyxl",
    "openpyxl.styles",
    "openpyxl.utils",
    "openpyxl.chart",
    "openpyxl.drawing",
    "openpyxl.drawing.image",
    "docx",
    "docx.shared",
    "docx.enum",
    "docx.enum.text",
    "docx.enum.table",
    "docx.enum.style",
    "docx.oxml",
    "pptx",
    "pptx.util",
    "pptx.dml.color",
    "pptx.enum.shapes",
    "pptx.enum.text",
    # 标准库（安全子集）
    "datetime",
    "math",
    "random",
    "string",
    "re",
    "json",
    "collections",
    "itertools",
    "functools",
    "decimal",
    "fractions",
    "statistics",
    "textwrap",
    "unicodedata",
    "copy",
}

# 禁止的代码模式（安全检查）
FORBIDDEN_PATTERNS = [
    r"\bexec\s*\(",
    r"\beval\s*\(",
    r"\bcompile\s*\(",
    r"\b__import__\s*\(",
    r"\bgetattr\s*\(",
    r"\bsetattr\s*\(",
    r"\bdelattr\s*\(",
    r"\bglobals\s*\(",
    r"\blocals\s*\(",
    r"\bvars\s*\(",
    r"\bdir\s*\(",
    r"\bopen\s*\(",  # 禁止直接 open，使用我们提供的安全版本
    r"\bos\.",  # 禁止 os 模块
    r"\bsys\.",  # 禁止 sys 模块
    r"\bsubprocess",
    r"\bshutil",
    r"\bpickle",
    r"\bsocket",
    r"\brequests",
    r"\burllib",
    r"\bhttp",
    r"\bftplib",
    r"\bsmtplib",
    r"\btelnetlib",
    r"\bparamiko",
    r"\bfabric",
    r"\bpexpect",
    r"\bpty",
    r"\btty",
    r"\btermios",
    r"\bctypes",
    r"\b_thread",
    r"\bthreading\.Thread",
    r"\bmultiprocessing",
    r"\bimportlib",
    r"\bbuiltins",
    r"\b__builtins__",
    r"\b__class__",
    r"\b__subclasses__",
    r"\b__mro__",
    r"\b__bases__",
]


class CodeExecutionError(Exception):
    """代码执行错误"""

    pass


class SecurityViolationError(Exception):
    """安全违规错误"""

    pass


class RestrictedCodeExecutor:
    """受限代码执行器"""

    def __init__(self, work_dir: Path, timeout: float = 30.0):
        """
        Args:
            work_dir: 工作目录，所有文件操作限制在此目录内
            timeout: 执行超时时间（秒）
        """
        self.work_dir = work_dir.resolve()
        self.timeout = timeout
        self._executor = ThreadPoolExecutor(max_workers=1)

    def _validate_code(self, code: str) -> tuple[bool, str]:
        """
        验证代码安全性

        Returns:
            (是否安全, 错误信息)
        """
        # 检查禁止的模式
        for pattern in FORBIDDEN_PATTERNS:
            if re.search(pattern, code):
                return False, f"检测到禁止的代码模式: {pattern}"

        # 检查 import 语句
        import_pattern = r"^\s*(?:from\s+(\S+)\s+import|import\s+(\S+))"
        for line in code.split("\n"):
            match = re.match(import_pattern, line)
            if match:
                module = match.group(1) or match.group(2)
                # 提取顶级模块名
                base_module = module.split(".")[0]
                full_module = module.split(" ")[0]  # 处理 "import x as y"

                # 检查是否在白名单中
                allowed = False
                for allowed_mod in ALLOWED_IMPORTS:
                    if (
                        full_module == allowed_mod
                        or full_module.startswith(allowed_mod + ".")
                        or allowed_mod.startswith(full_module + ".")
                    ):
                        allowed = True
                        break

                if not allowed and base_module not in {
                    "openpyxl",
                    "docx",
                    "pptx",
                    "datetime",
                    "math",
                    "random",
                    "string",
                    "re",
                    "json",
                    "collections",
                    "itertools",
                    "functools",
                    "decimal",
                    "fractions",
                    "statistics",
                    "textwrap",
                    "unicodedata",
                    "copy",
                }:
                    return False, f"禁止导入模块: {module}"

        return True, ""

    def _create_safe_globals(self, output_path: Path) -> dict[str, Any]:
        """创建安全的全局命名空间"""
        import copy
        import datetime
        import json
        import math
        import random
        import re as re_module
        import string

        # 安全的文件保存函数
        def safe_save(obj: Any, filename: str) -> str:
            """安全地保存文件到工作目录"""
            # 验证文件名
            safe_name = Path(filename).name  # 只取文件名，防止路径穿越
            if not safe_name:
                raise ValueError("无效的文件名")

            file_path = output_path / safe_name

            # 确保路径在工作目录内
            if not file_path.resolve().is_relative_to(output_path.resolve()):
                raise SecurityViolationError("禁止访问工作目录外的文件")

            # 根据对象类型保存
            if hasattr(obj, "save"):
                obj.save(str(file_path))
            else:
                raise TypeError(f"对象类型 {type(obj)} 不支持保存")

            return str(file_path)

        # 构建安全的全局命名空间
        safe_globals = {
            "__builtins__": {
                # 安全的内置函数
                "abs": abs,
                "all": all,
                "any": any,
                "bin": bin,
                "bool": bool,
                "chr": chr,
                "dict": dict,
                "divmod": divmod,
                "enumerate": enumerate,
                "filter": filter,
                "float": float,
                "format": format,
                "frozenset": frozenset,
                "hash": hash,
                "hex": hex,
                "int": int,
                "isinstance": isinstance,
                "issubclass": issubclass,
                "iter": iter,
                "len": len,
                "list": list,
                "map": map,
                "max": max,
                "min": min,
                "next": next,
                "oct": oct,
                "ord": ord,
                "pow": pow,
                "print": print,  # 允许 print 用于调试
                "range": range,
                "repr": repr,
                "reversed": reversed,
                "round": round,
                "set": set,
                "slice": slice,
                "sorted": sorted,
                "str": str,
                "sum": sum,
                "tuple": tuple,
                "type": type,
                "zip": zip,
                # 异常类
                "Exception": Exception,
                "ValueError": ValueError,
                "TypeError": TypeError,
                "KeyError": KeyError,
                "IndexError": IndexError,
                "AttributeError": AttributeError,
                "RuntimeError": RuntimeError,
                "StopIteration": StopIteration,
                # 其他
                "True": True,
                "False": False,
                "None": None,
            },
            # 预导入的安全模块
            "datetime": datetime,
            "math": math,
            "random": random,
            "string": string,
            "re": re_module,
            "json": json,
            "copy": copy,
            # 工具函数
            "save_file": safe_save,
            "OUTPUT_DIR": str(output_path),
        }

        return safe_globals

    def _execute_sync(self, code: str, output_path: Path) -> dict[str, Any]:
        """同步执行代码"""
        # 验证代码安全性
        is_safe, error = self._validate_code(code)
        if not is_safe:
            raise SecurityViolationError(error)

        # 创建安全的执行环境
        safe_globals = self._create_safe_globals(output_path)
        safe_locals: dict[str, Any] = {}

        try:
            # 编译代码
            compiled = compile(code, "<llm_generated>", "exec")

            # 执行代码
            exec(compiled, safe_globals, safe_locals)

            # 返回结果
            return {
                "success": True,
                "locals": safe_locals,
                "output_dir": str(output_path),
            }

        except SecurityViolationError:
            raise
        except Exception as e:
            tb = traceback.format_exc()
            raise CodeExecutionError(f"代码执行失败: {e}\n{tb}")

    async def execute(self, code: str) -> dict[str, Any]:
        """
        异步执行 LLM 生成的代码

        Args:
            code: Python 代码字符串

        Returns:
            执行结果字典，包含生成的文件路径等信息
        """
        # 创建输出目录
        output_path = self.work_dir
        output_path.mkdir(parents=True, exist_ok=True)

        loop = asyncio.get_event_loop()

        try:
            # 在线程池中执行，带超时
            result = await asyncio.wait_for(
                loop.run_in_executor(
                    self._executor, self._execute_sync, code, output_path
                ),
                timeout=self.timeout,
            )
            return result

        except asyncio.TimeoutError:
            raise CodeExecutionError(f"代码执行超时（{self.timeout}秒）")

    def cleanup(self):
        """清理资源"""
        self._executor.shutdown(wait=False)


# 代码生成提示模板
CODE_GENERATION_PROMPT = """你是一个专业的 Python 代码生成器，专门生成创建 Office 文档的代码。

## 可用库
- openpyxl: Excel 文件（.xlsx）
- python-docx: Word 文件（.docx）
- python-pptx: PowerPoint 文件（.pptx）
- datetime, math, random, string, re, json: 标准库

## 重要规则
1. 必须使用 `save_file(obj, filename)` 函数保存文件
2. 文件会保存到 OUTPUT_DIR 目录
3. 不要使用 open()、os、sys 等禁止的函数
4. 代码必须是完整可执行的

## Excel 示例
```python
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

wb = Workbook()
ws = wb.active
ws.title = "数据表"

# 设置表头样式
header_font = Font(bold=True, color="FFFFFF")
header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")

headers = ["姓名", "年龄", "城市", "职业"]
for col, header in enumerate(headers, 1):
    cell = ws.cell(row=1, column=col, value=header)
    cell.font = header_font
    cell.fill = header_fill
    cell.alignment = Alignment(horizontal="center")

# 添加数据
data = [
    ["张三", 28, "北京", "工程师"],
    ["李四", 32, "上海", "设计师"],
]
for row_idx, row_data in enumerate(data, 2):
    for col_idx, value in enumerate(row_data, 1):
        ws.cell(row=row_idx, column=col_idx, value=value)

# 调整列宽
for col in ws.columns:
    ws.column_dimensions[col[0].column_letter].width = 15

save_file(wb, "数据表.xlsx")
```

## Word 示例
```python
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

doc = Document()

# 添加标题
title = doc.add_heading("项目报告", 0)
title.alignment = WD_ALIGN_PARAGRAPH.CENTER

# 添加段落
para = doc.add_paragraph()
run = para.add_run("这是报告的正文内容。")
run.font.size = Pt(12)

# 添加表格
table = doc.add_table(rows=3, cols=3)
table.style = "Table Grid"
for i, row in enumerate(table.rows):
    for j, cell in enumerate(row.cells):
        cell.text = f"单元格 {i},{j}"

save_file(doc, "报告.docx")
```

## PPT 示例
```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# 添加标题页
slide_layout = prs.slide_layouts[0]  # 标题布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "演示文稿标题"
subtitle = slide.placeholders[1]
subtitle.text = "副标题内容"

# 添加内容页
slide_layout = prs.slide_layouts[1]  # 标题和内容布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "第一章"
body = slide.placeholders[1]
tf = body.text_frame
tf.text = "要点一"
p = tf.add_paragraph()
p.text = "要点二"
p.level = 0

save_file(prs, "演示文稿.pptx")
```

现在请根据用户需求生成代码：
"""


def extract_code_from_response(response: str) -> str:
    """从 LLM 响应中提取代码块"""
    # 尝试匹配 ```python ... ``` 代码块
    pattern = r"```(?:python)?\s*\n(.*?)```"
    matches = re.findall(pattern, response, re.DOTALL)

    if matches:
        # 返回最后一个代码块（通常是完整代码）
        return matches[-1].strip()

    # 如果没有代码块，假设整个响应就是代码
    return response.strip()
