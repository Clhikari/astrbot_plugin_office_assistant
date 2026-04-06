import asyncio
import importlib.util
import json
import re
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from pathlib import Path

from pydantic import TypeAdapter, ValidationError

from astrbot.api import logger
from astrbot.api.event import AstrMessageEvent
from astrbot.core.message.message_event_result import MessageChain

from .constants import OFFICE_EXTENSIONS, OFFICE_LIBS, OfficeType
from ._executor_mixin import ExecutorOwnerMixin
from .document_core.models.blocks import (
    DocumentBlock,
    HeadingBlock,
    ParagraphBlock,
    TableBlock,
)
from .document_core.models.document import (
    DocumentMetadata,
    DocumentModel,
    DocumentStatus,
)
from .domain.document.contracts import (
    AddBlocksRequest,
    BlockInput,
    CreateDocumentRequest,
    normalize_raw_block_payloads,
)
from .domain.document.render_backends import (
    DocumentRenderBackendConfig,
    build_document_render_backends,
    render_document_with_backends,
)
from .domain.document.session_store import DocumentSessionStore

_BLOCK_INPUT_ADAPTER = TypeAdapter(BlockInput)


class OfficeGenerator(ExecutorOwnerMixin):
    """Office文件生成器"""

    def __init__(
        self,
        data_path: Path,
        executor: ThreadPoolExecutor | None = None,
        render_backend_config: DocumentRenderBackendConfig | None = None,
    ):
        self.data_path = data_path
        self.support = self._check_support()
        self._render_backend_config = render_backend_config
        self._init_executor(executor, label="文件生成器")

    # 定义映射表
    _GENERATORS = {
        OfficeType.WORD: "_generate_word",
        OfficeType.EXCEL: "_generate_excel",
        OfficeType.POWERPOINT: "_generate_powerpoint",
    }

    async def _generate_word(self, file_path: Path, content: dict):
        """异步生成 Word"""
        loop = asyncio.get_running_loop()
        await loop.run_in_executor(
            self._require_executor(), self._generate_word_sync, file_path, content
        )

    async def _generate_excel(self, file_path: Path, content: dict):
        """异步生成 Excel"""
        loop = asyncio.get_running_loop()
        await loop.run_in_executor(
            self._require_executor(), self._generate_excel_sync, file_path, content
        )

    async def _generate_powerpoint(self, file_path: Path, content: dict):
        """异步生成 PPT"""
        loop = asyncio.get_running_loop()
        await loop.run_in_executor(
            self._require_executor(), self._generate_ppt_sync, file_path, content
        )

    def _check_support(self) -> dict[OfficeType, bool]:
        """检查Office库支持"""
        support = {}

        for office_type in OfficeType:
            module_name, package_name = OFFICE_LIBS[office_type]
            available = importlib.util.find_spec(module_name) is not None
            support[office_type] = available

            if not available:
                logger.warning(f"[文件生成器] {package_name} 未安装")

        return support

    async def generate(
        self,
        event: AstrMessageEvent,
        office_type: OfficeType,
        filename: str,
        content: dict,
    ):
        """生成Office文件"""
        if not self.support.get(office_type, False):
            await event.send(
                MessageChain().message(
                    f"[文件生成器] {office_type}文件生成不支持，缺少相关库"
                )
            )
            return ""

        try:
            resolved_filename = (
                content.get("filename") or filename or f"{office_type}_file"
            )
            payload = content.get("content", {})

            # 解析content
            if isinstance(payload, str):
                try:
                    payload = json.loads(payload)
                except json.JSONDecodeError:
                    payload = self._create_default_content(office_type, payload)

            # 清理文件名并添加扩展名
            resolved_filename = self._sanitize_filename(resolved_filename)
            extension = OFFICE_EXTENSIONS[office_type]

            if not resolved_filename.endswith(extension):
                resolved_filename = resolved_filename + extension

            file_path = self._get_unique_filepath(resolved_filename)
            generator = getattr(self, self._GENERATORS[office_type])
            await generator(file_path, payload)

            logger.info(f"[文件生成器] Office文件已生成: {file_path}")
            return file_path

        except Exception as e:
            logger.error(f"[文件生成器] 生成Office文件失败: {e}", exc_info=True)
            return ""

    def _create_default_content(self, file_type: OfficeType, text: str) -> dict:
        """创建默认的内容结构，智能解析文本格式"""
        if file_type == OfficeType.WORD:
            return self._parse_word_content(text)
        elif file_type == OfficeType.EXCEL:
            return self._parse_excel_content(text)
        elif file_type == OfficeType.POWERPOINT:
            return self._parse_ppt_content(text)
        return {}

    def _parse_word_content(self, text: str) -> dict:
        """解析 Word 内容格式

        支持：
        - 空行分隔段落
        - 第一行作为标题（如果后面有空行）
        """
        text = text.strip()
        if not text:
            return {"paragraphs": ["（空文档）"]}

        # 按空行分割段落
        paragraphs = []
        current_para = []

        for line in text.split("\n"):
            if line.strip():
                current_para.append(line.strip())
            else:
                if current_para:
                    paragraphs.append(" ".join(current_para))
                    current_para = []

        if current_para:
            paragraphs.append(" ".join(current_para))

        if not paragraphs:
            return {"paragraphs": [text]}

        # 如果只有一个段落，直接返回
        if len(paragraphs) == 1:
            return {"paragraphs": paragraphs}

        # 第一段作为标题
        title = paragraphs[0]
        body_paragraphs = paragraphs[1:]

        return {"title": title, "paragraphs": body_paragraphs}

    def _parse_excel_content(self, text: str) -> dict:
        """解析 Excel 内容格式

        支持：
        - 用 | 分隔每个单元格
        - 换行分隔每一行
        """
        lines = [line.strip() for line in text.strip().split("\n") if line.strip()]
        if not lines:
            return {"sheets": [{"name": "Sheet1", "data": [["（空表格）"]]}]}

        data = []
        for line in lines:
            if "|" in line:
                cells = [cell.strip() for cell in line.split("|")]
            else:
                cells = [line]
            data.append(cells)

        return {"sheets": [{"name": "Sheet1", "data": data}]}

    def _parse_ppt_content(self, text: str) -> dict:
        """解析 PPT 内容格式

        支持：
        - 用 [幻灯片 N] 或 [Slide N] 标记每页
        - 标记后第一行作为标题
        - 后续行作为内容要点
        """

        text = text.strip()
        if not text:
            return {"slides": [{"title": "空白幻灯片", "content": []}]}

        # 匹配 [幻灯片 N]、[幻灯片N]、[Slide N]、[slide N] 等格式
        slide_pattern = r"\[(?:幻灯片|Slide|slide)\s*\d+\]"

        # 检查是否有幻灯片标记
        if not re.search(slide_pattern, text):
            # 没有标记，尝试按空行分割成多页
            return self._parse_ppt_by_blank_lines(text)

        # 按幻灯片标记分割
        slides = []
        parts = re.split(f"({slide_pattern})", text)

        current_slide = None
        for part in parts:
            part = part.strip()
            if not part:
                continue

            if re.match(slide_pattern, part):
                # 这是一个幻灯片标记，准备新幻灯片
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": "", "content": []}
            elif current_slide is not None:
                # 解析幻灯片内容
                lines = [ln.strip() for ln in part.split("\n") if ln.strip()]
                if lines:
                    current_slide["title"] = lines[0]
                    current_slide["content"] = lines[1:] if len(lines) > 1 else []

        if current_slide:
            slides.append(current_slide)

        if not slides:
            return {"slides": [{"title": "内容", "content": [text]}]}

        return {"slides": slides}

    def _parse_ppt_by_blank_lines(self, text: str) -> dict:
        """按空行分割 PPT 内容为多页"""
        slides = []
        current_lines = []

        for line in text.split("\n"):
            if line.strip():
                current_lines.append(line.strip())
            else:
                if current_lines:
                    title = current_lines[0]
                    content = current_lines[1:] if len(current_lines) > 1 else []
                    slides.append({"title": title, "content": content})
                    current_lines = []

        if current_lines:
            title = current_lines[0]
            content = current_lines[1:] if len(current_lines) > 1 else []
            slides.append({"title": title, "content": content})

        if not slides:
            return {"slides": [{"title": "内容", "content": [text]}]}

        return {"slides": slides}

    def _generate_word_sync(self, file_path: Path, content: dict):
        """生成Word文档"""
        document_model = self._build_document_model(file_path, content)
        render_document_with_backends(
            document_model,
            file_path,
            build_document_render_backends(
                "word",
                self._render_backend_config,
            ),
        )

    def _build_document_model(
        self, file_path: Path, content: dict | DocumentModel
    ) -> DocumentModel:
        if isinstance(content, DocumentModel):
            document_model = content.model_copy(deep=True)
            document_model.metadata.preferred_filename = file_path.name
            return document_model

        if self._looks_like_document_model_payload(content):
            try:
                metadata_data = content.get("metadata", {})
                metadata = (
                    metadata_data
                    if isinstance(metadata_data, DocumentMetadata)
                    else DocumentMetadata.model_validate(metadata_data)
                )
                return self._build_structured_document_model(
                    file_path=file_path,
                    content=content,
                    metadata=metadata,
                )
            except ValidationError as exc:
                logger.warning(
                    f"[文件生成器] 文档块模型校验失败，回退到旧版 Word 内容适配: {exc}"
                )

        return self._build_legacy_document_model(file_path, content)

    def _build_word_document_model(
        self, file_path: Path, content: dict | DocumentModel
    ) -> DocumentModel:
        return self._build_document_model(file_path, content)

    def _looks_like_document_model_payload(self, content: dict) -> bool:
        return isinstance(content, dict) and "blocks" in content

    def _build_structured_document_model(
        self,
        *,
        file_path: Path,
        content: dict,
        metadata: DocumentMetadata,
    ) -> DocumentModel:
        metadata = metadata.model_copy(deep=True)
        metadata.preferred_filename = file_path.name

        temp_store = DocumentSessionStore(workspace_dir=self.data_path)
        created_document = temp_store.create_document(
            CreateDocumentRequest(
                session_id=str(content.get("session_id") or ""),
                title=metadata.title,
                output_name=file_path.name,
                theme_name=metadata.theme_name,
                table_template=metadata.table_template,
                density=metadata.density,
                accent_color=metadata.accent_color,
                header_footer=metadata.header_footer.model_dump(
                    mode="json",
                    exclude_none=True,
                ),
                document_style=metadata.document_style.model_dump(
                    mode="json",
                    exclude_none=True,
                ),
            )
        )

        normalized_payloads = normalize_raw_block_payloads(list(content.get("blocks") or []))
        validated_blocks: list[BlockInput] = []
        for idx, block_payload in enumerate(normalized_payloads):
            try:
                validated_blocks.append(_BLOCK_INPUT_ADAPTER.validate_python(block_payload))
            except ValidationError as exc:
                logger.warning(
                    "[文件生成器] 跳过无效文档块 index=%s file=%s: %s",
                    idx,
                    file_path,
                    exc,
                )

        if validated_blocks:
            temp_store.add_blocks(
                AddBlocksRequest(
                    document_id=created_document.document_id,
                    blocks=validated_blocks,
                )
            )

        status_value = content.get("status", DocumentStatus.DRAFT)
        if isinstance(status_value, str):
            try:
                status_value = DocumentStatus(status_value)
            except ValueError:
                status_value = DocumentStatus.DRAFT

        created_document.document_id = str(content.get("document_id") or file_path.stem)
        created_document.session_id = str(content.get("session_id") or "")
        created_document.status = status_value
        created_document.output_path = str(file_path)
        created_document.metadata.preferred_filename = file_path.name
        return created_document

    def _build_structured_word_document_model(
        self,
        *,
        file_path: Path,
        content: dict,
        metadata: DocumentMetadata,
    ) -> DocumentModel:
        return self._build_structured_document_model(
            file_path=file_path,
            content=content,
            metadata=metadata,
        )

    def _build_legacy_document_model(
        self, file_path: Path, content: dict
    ) -> DocumentModel:
        metadata = DocumentMetadata(
            title=str(content.get("title", "")),
            preferred_filename=file_path.name,
        )
        blocks: list[DocumentBlock] = []

        paragraphs = content.get("paragraphs", [])
        if isinstance(paragraphs, str):
            paragraphs = [paragraphs]

        for para_text in paragraphs:
            normalized = str(para_text).strip()
            if normalized:
                blocks.append(ParagraphBlock(text=normalized))

        table_block = self._build_legacy_table_block(content.get("table"))
        if table_block is not None:
            blocks.append(table_block)

        return DocumentModel(
            document_id=file_path.stem,
            format="word",
            metadata=metadata,
            blocks=blocks,
            output_path=str(file_path),
        )

    def _build_legacy_word_document_model(
        self, file_path: Path, content: dict
    ) -> DocumentModel:
        return self._build_legacy_document_model(file_path, content)

    def _build_legacy_table_block(self, table_data: object) -> TableBlock | None:
        if not isinstance(table_data, list) or not table_data:
            return None

        normalized_rows: list[list[str]] = []
        for row in table_data:
            if isinstance(row, list):
                normalized_rows.append(
                    ["" if cell is None else str(cell) for cell in row]
                )
            else:
                normalized_rows.append([str(row)])

        if not normalized_rows:
            return None

        headers = normalized_rows[0]
        body_rows = normalized_rows[1:] if len(normalized_rows) > 1 else []
        return TableBlock(headers=headers, rows=body_rows)

    def _generate_excel_sync(self, file_path: Path, content: dict):
        """生成Excel表格"""
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill

        wb = Workbook()

        # 删除默认sheet
        if "Sheet" in wb.sheetnames:
            wb.remove(wb["Sheet"])

        # 添加工作表
        sheets = content.get("sheets", [])
        if not sheets:
            sheets = [{"name": "Sheet1", "data": [["无数据"]]}]

        for sheet_info in sheets:
            sheet_name = sheet_info.get("name", "Sheet1")
            sheet_data = sheet_info.get("data", [])

            ws = wb.create_sheet(title=sheet_name)

            # 写入数据
            for row_idx, row_data in enumerate(sheet_data, 1):
                for col_idx, cell_value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)

                    # 第一行加粗
                    if row_idx == 1:
                        cell.font = Font(bold=True)
                        cell.fill = PatternFill(
                            start_color="DDDDDD", end_color="DDDDDD", fill_type="solid"
                        )

                    cell.alignment = Alignment(horizontal="left", vertical="center")

        wb.save(str(file_path))

    def _generate_ppt_sync(self, file_path: Path, content: dict):
        """生成PowerPoint演示文稿"""
        from pptx import Presentation

        prs = Presentation()

        # 添加幻灯片
        slides_data = content.get("slides", [])
        if not slides_data:
            slides_data = [{"title": "标题", "content": ["内容"]}]

        for slide_info in slides_data:
            # 使用标题和内容布局
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            title = slide_info.get("title", "")
            if title:
                slide.shapes.title.text = title

            # 添加内容
            content_list = slide_info.get("content", [])
            if isinstance(content_list, str):
                content_list = [content_list]

            if content_list and len(slide.shapes) > 1:
                text_frame = slide.shapes[1].text_frame
                text_frame.clear()

                for item in content_list:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 0

        prs.save(str(file_path))

    def _sanitize_filename(self, filename: str) -> str:
        """
        清理文件名

        - 移除非法字符，仅保留字母、数字、中文、空格、连字符、下划线和点号
        - 移除开头的点号（避免创建隐藏文件）
        - 合并多个连续的点号
        - 限制文件名长度
        """

        # 保留字母、数字、中文、空格、连字符、下划线、点号
        filename = "".join(
            c for c in filename if c.isalnum() or c in (" ", "-", "_", ".")
        ).strip()

        # 移除开头的点号
        filename = filename.lstrip(".")

        # 合并多个连续的点号为单个点号
        filename = re.sub(r"\.{2,}", ".", filename)

        # 移除末尾多余的点号（扩展名前的点号除外）
        # 先分离扩展名
        if "." in filename:
            name_part, ext_part = filename.rsplit(".", 1)
            name_part = name_part.rstrip(".")
            filename = f"{name_part}.{ext_part}" if name_part else ext_part
        else:
            filename = filename.rstrip(".")

        # 限制长度（保留扩展名）
        max_length = 200
        if len(filename) > max_length:
            if "." in filename:
                name_part, ext_part = filename.rsplit(".", 1)
                name_part = name_part[: max_length - len(ext_part) - 1]
                filename = f"{name_part}.{ext_part}"
            else:
                filename = filename[:max_length]

        # 如果清理后为空，使用默认文件名
        if not filename or filename == ".":
            filename = f"office_file_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

        return filename

    def _get_unique_filepath(self, filename: str) -> Path:
        """获取文件路径（覆盖已存在的同名文件）"""
        return self.data_path / filename

    def cleanup(self):
        """清理资源"""
        self._shutdown_executor()
