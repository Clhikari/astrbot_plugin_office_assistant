import hashlib
import re
import subprocess
from collections.abc import Generator
from contextlib import contextmanager, suppress
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

from astrbot.api import logger

from .compat import _ANTIWORD_AVAILABLE, _IS_WINDOWS, _WIN32COM_AVAILABLE

if _WIN32COM_AVAILABLE:
    import pythoncom
    import win32com.client


@contextmanager
def com_application(app_name: str) -> Generator:
    """Windows COM 应用上下文管理器

    自动处理 COM 初始化/反初始化和应用退出，避免资源泄漏。

    Args:
        app_name: COM 应用名称，如 "Word.Application"

    Yields:
        COM 应用对象

    Example:
        with com_application("Word.Application") as app:
            doc = app.Documents.Open(path)
            # ...
    """
    if not _WIN32COM_AVAILABLE:
        raise RuntimeError("win32com 不可用，请安装 pywin32: pip install pywin32")

    app = None
    try:
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch(app_name)
        app.Visible = False
        if hasattr(app, "DisplayAlerts"):
            app.DisplayAlerts = False
        yield app
    finally:
        if app:
            with suppress(Exception):
                app.Quit()
        with suppress(Exception):
            pythoncom.CoUninitialize()


def format_file_size(size: int | float) -> str:
    """格式化文件大小"""
    if size <= 0:
        return "0 B"
    for unit in ("B", "KB", "MB", "GB"):
        if size < 1024:
            return f"{size:.2f} {unit}"
        size /= 1024
    return f"{size:.2f} TB"


def safe_error_message(error: Exception, context: str = "") -> str:
    """
    生成安全的错误消息，隐藏敏感路径信息

    Args:
        error: 异常对象
        context: 错误上下文描述

    Returns:
        不包含敏感信息的错误消息
    """
    error_str = str(error)

    # 移除可能包含路径的信息
    # 常见模式: 'D:\\path\\to\\file' 或 '/path/to/file'

    # Windows 路径模式
    error_str = re.sub(r"[A-Za-z]:\\[^\s'\"]+", "[路径已隐藏]", error_str)
    # Unix 路径模式
    error_str = re.sub(
        r"/(?:home|tmp|var|usr|opt|data)[^\s'\"]*", "[路径已隐藏]", error_str
    )

    if context:
        return f"{context}: {error_str}"
    return error_str


WordItemType = Literal["text", "image"]
WORD_ITEM_TEXT: WordItemType = "text"
WORD_ITEM_IMAGE: WordItemType = "image"


@dataclass(slots=True)
class ExtractedWordItem:
    type: WordItemType
    text: str | None = None
    image_path: Path | None = None
    image_index: int | None = None


@dataclass(slots=True)
class ExtractedWordContent:
    text: str | None = None
    image_paths: list[Path] = field(default_factory=list)
    items: list[ExtractedWordItem] = field(default_factory=list)
    image_count: int = 0


@dataclass(slots=True)
class ExtractedExcelSheet:
    name: str
    text: str


_EXCEL_PREVIEW_MAX_ROWS_PER_SHEET = 200
_EXCEL_PREVIEW_MAX_CHARS_PER_SHEET = 12_000
_EXCEL_PREVIEW_MAX_TOTAL_CHARS = 24_000
_EXCEL_PREVIEW_NOTE_RESERVE = 120


def _build_excel_sheet_preview(
    *,
    sheet_name: str,
    rows,
    remaining_chars: int,
) -> tuple[ExtractedExcelSheet, int, bool]:
    sheet_limit = min(_EXCEL_PREVIEW_MAX_CHARS_PER_SHEET, remaining_chars)
    content_limit = max(sheet_limit - _EXCEL_PREVIEW_NOTE_RESERVE, 0)
    lines: list[str] = []
    content_chars = 0
    shown_rows = 0
    row_limit_hit = False
    char_limit_hit = False
    workbook_limit_hit = False

    for row in rows:
        if shown_rows >= _EXCEL_PREVIEW_MAX_ROWS_PER_SHEET:
            row_limit_hit = True
            break

        row_text = "\t".join(_normalize_excel_preview_value(value) for value in row)
        line_cost = len(row_text) + (1 if lines else 0)
        if content_chars + line_cost <= content_limit:
            lines.append(row_text)
            content_chars += line_cost
            shown_rows += 1
            continue

        available_chars = content_limit - content_chars - (1 if lines else 0)
        if available_chars > 0:
            truncated_row = row_text[:available_chars].rstrip()
            if truncated_row:
                lines.append(truncated_row)
                content_chars += len(truncated_row) + (1 if len(lines) > 1 else 0)
                shown_rows += 1

        char_limit_hit = True
        workbook_limit_hit = sheet_limit < _EXCEL_PREVIEW_MAX_CHARS_PER_SHEET
        break

    note_parts: list[str] = []
    if row_limit_hit:
        note_parts.append(f"仅展示前 {shown_rows} 行")
    if workbook_limit_hit:
        note_parts.append("工作簿总预览已达到上限")
    elif char_limit_hit:
        note_parts.append("单个 Sheet 预览已达到上限")
    if note_parts:
        lines.append(f"[已截断：{'；'.join(note_parts)}]")

    text = "\n".join(lines)
    return (
        ExtractedExcelSheet(name=sheet_name, text=text),
        len(text),
        workbook_limit_hit,
    )


def _normalize_excel_preview_value(value: object) -> str:
    if value is None:
        return ""
    return re.sub(r"[\r\n\t]+", " ", str(value))


def format_extracted_word_content(
    content: ExtractedWordContent | None,
    *,
    workspace_root: Path | None = None,
    include_image_paths: bool = False,
    item_separator: str = "\n",
) -> str | None:
    if content is None:
        return None

    normalized_root = workspace_root.resolve() if workspace_root is not None else None

    def build_image_line(index: int | None, image_path: Path | None) -> str | None:
        label = f"[插图{index}]" if index is not None else "[插图]"
        if not include_image_paths:
            return label
        if image_path is None:
            return None
        if normalized_root is not None and image_path.is_relative_to(normalized_root):
            display_path = image_path.relative_to(normalized_root).as_posix()
        else:
            display_path = image_path.name
        return f"{label} {display_path}"

    if content.items:
        item_lines: list[str] = []
        for item in content.items:
            if item.type == WORD_ITEM_TEXT:
                text = (item.text or "").strip()
                if text:
                    item_lines.append(text)
                continue
            if item.type == WORD_ITEM_IMAGE:
                image_line = build_image_line(item.image_index, item.image_path)
                if image_line:
                    item_lines.append(image_line)
        if item_lines:
            return item_separator.join(item_lines)

    parts: list[str] = []
    if content.text:
        parts.append(content.text)

    if content.image_paths:
        for index, image_path in enumerate(content.image_paths, start=1):
            image_line = build_image_line(index, image_path)
            if image_line:
                parts.append(image_line)

    return item_separator.join(parts) if parts else None


def extract_word_content(
    file_path: Path,
    workspace_root: Path | None = None,
    *,
    include_images: bool = True,
) -> ExtractedWordContent | None:
    """Extract structured Word content for downstream tool formatting."""
    suffix = file_path.suffix.lower()

    if suffix == ".doc":
        if _ANTIWORD_AVAILABLE:
            text = _extract_doc_text_antiword(file_path)
        elif _WIN32COM_AVAILABLE:
            text = _extract_doc_text_win32com(file_path)
        else:
            if _IS_WINDOWS:
                logger.warning(
                    "读取 .doc 文件需要 pywin32 和 Microsoft Word，请安装: pip install pywin32"
                )
            else:
                logger.warning(
                    "读取 .doc 文件需要 antiword，请安装: apt install antiword (Linux) 或 brew install antiword (macOS)"
                )
            return None
        if not text:
            return None
        return ExtractedWordContent(text=text)

    try:
        from docx import Document

        doc = Document(file_path)
        image_parts = _collect_docx_image_parts(doc)
        image_count = len(image_parts)
        image_rel_paths = _extract_docx_images(
            image_parts,
            file_path,
            workspace_root,
            include_images=include_images,
        )
        items = _extract_docx_items(
            doc,
            image_rel_paths,
            include_images=include_images,
        )
        if not items and image_count == 0:
            return None
        image_paths = [
            item.image_path
            for item in items
            if item.type == WORD_ITEM_IMAGE and item.image_path is not None
        ]
        text_items = [
            item.text.strip()
            for item in items
            if item.type == WORD_ITEM_TEXT and item.text and item.text.strip()
        ]
        text = "\n".join(text_items) if text_items else None
        return ExtractedWordContent(
            text=text,
            image_paths=image_paths,
            items=items,
            image_count=image_count,
        )
    except ImportError:
        return None
    except Exception as e:
        logger.warning(f"Word 文本提取失败: {e}", exc_info=True)
        return None


def extract_word_text(
    file_path: Path,
    workspace_root: Path | None = None,
) -> str | None:
    """提取 Word 文档文本（支持 .docx 和 .doc）"""
    extracted = extract_word_content(file_path, workspace_root)
    return format_extracted_word_content(
        extracted,
        workspace_root=workspace_root,
        include_image_paths=True,
        item_separator="\n\n",
    )


def _collect_docx_image_parts(doc) -> list[tuple[str, bytes, str]]:
    image_parts = []
    for rel_id, rel in doc.part.rels.items():
        if "image" not in rel.reltype:
            continue
        target_part = getattr(rel, "target_part", None)
        if target_part is None:
            continue
        blob = getattr(target_part, "blob", None)
        partname = getattr(target_part, "partname", "")
        if not blob:
            continue
        image_parts.append((rel_id, blob, Path(str(partname)).suffix or ".bin"))

    return image_parts


def _extract_docx_images(
    image_parts: list[tuple[str, bytes, str]],
    file_path: Path,
    workspace_root: Path | None,
    *,
    include_images: bool,
) -> dict[str, Path]:
    if workspace_root is None or not include_images or not image_parts:
        return {}

    workspace_root = workspace_root.resolve()
    image_dir = _build_docx_asset_dir(file_path, workspace_root)
    image_dir.mkdir(parents=True, exist_ok=True)

    digest_to_path: dict[str, Path] = {}
    rel_paths: dict[str, Path] = {}
    image_index = 1
    for rel_id, blob, suffix in image_parts:
        digest = hashlib.md5(blob).hexdigest()
        output_path = digest_to_path.get(digest)
        if output_path is None:
            output_path = image_dir / f"image_{image_index:02d}{suffix}"
            output_path.write_bytes(blob)
            digest_to_path[digest] = output_path
            image_index += 1
        rel_paths[rel_id] = output_path
    return rel_paths


def _extract_docx_items(
    doc,
    image_rel_paths: dict[str, Path],
    *,
    include_images: bool,
) -> list[ExtractedWordItem]:
    items: list[ExtractedWordItem] = []

    for body_child in doc.element.body.iterchildren():
        local_name = body_child.tag.rsplit("}", 1)[-1]
        if local_name != "p":
            continue

        paragraph_buffer: list[str] = []
        for child in body_child.iterchildren():
            child_name = child.tag.rsplit("}", 1)[-1]
            if child_name in {"r", "hyperlink", "smartTag", "sdt", "ins"}:
                _collect_paragraph_items(
                    child,
                    image_rel_paths=image_rel_paths,
                    paragraph_buffer=paragraph_buffer,
                    items=items,
                    include_images=include_images,
                )

        paragraph_text = "".join(paragraph_buffer).strip()
        if paragraph_text:
            items.append(ExtractedWordItem(type=WORD_ITEM_TEXT, text=paragraph_text))

    image_index = 1
    for item in items:
        if item.type == WORD_ITEM_IMAGE:
            item.image_index = image_index
            image_index += 1
    return items


def _collect_paragraph_items(
    element,
    *,
    image_rel_paths: dict[str, Path],
    paragraph_buffer: list[str],
    items: list[ExtractedWordItem],
    include_images: bool,
) -> None:
    local_name = element.tag.rsplit("}", 1)[-1]

    if local_name == "r":
        for child in element.iterchildren():
            child_name = child.tag.rsplit("}", 1)[-1]
            if child_name == "t":
                if child.text:
                    paragraph_buffer.append(child.text)
            elif child_name == "tab":
                paragraph_buffer.append("\t")
            elif child_name in {"br", "cr"}:
                paragraph_buffer.append("\n")
            elif child_name in {"drawing", "object", "pict"}:
                paragraph_text = "".join(paragraph_buffer).strip()
                if paragraph_text:
                    items.append(
                        ExtractedWordItem(type=WORD_ITEM_TEXT, text=paragraph_text)
                    )
                paragraph_buffer.clear()
                if include_images:
                    for image_path in _extract_embedded_image_paths(
                        child, image_rel_paths
                    ):
                        items.append(
                            ExtractedWordItem(
                                type=WORD_ITEM_IMAGE,
                                image_path=image_path,
                            )
                        )
        return

    for child in element.iterchildren():
        child_name = child.tag.rsplit("}", 1)[-1]
        if child_name in {"r", "hyperlink", "smartTag", "sdt", "ins"}:
            _collect_paragraph_items(
                child,
                image_rel_paths=image_rel_paths,
                paragraph_buffer=paragraph_buffer,
                items=items,
                include_images=include_images,
            )


def _extract_embedded_image_paths(
    element, image_rel_paths: dict[str, Path]
) -> list[Path]:
    image_paths: list[Path] = []
    seen_rel_ids: set[str] = set()

    for node in element.iter():
        for attr_name, attr_value in node.attrib.items():
            local_name = attr_name.rsplit("}", 1)[-1]
            if local_name not in {"embed", "id"}:
                continue
            image_path = image_rel_paths.get(attr_value)
            if image_path is None or attr_value in seen_rel_ids:
                continue
            seen_rel_ids.add(attr_value)
            image_paths.append(image_path)

    return image_paths


def _build_docx_asset_dir(file_path: Path, workspace_root: Path) -> Path:
    digest = hashlib.md5(str(file_path.resolve()).encode("utf-8")).hexdigest()[:8]
    return workspace_root / ".read_assets" / f"{file_path.stem}_{digest}"


def _extract_doc_text_antiword(file_path: Path) -> str | None:
    """使用 antiword 提取 .doc 文件文本（Linux/macOS）"""
    try:
        result = subprocess.run(
            ["antiword", str(file_path.resolve())],
            capture_output=True,
            text=True,
            errors="replace",
            timeout=30,
        )
        if result.returncode == 0:
            return result.stdout.strip() if result.stdout else None
        else:
            logger.warning(f"antiword 执行失败: {result.stderr}")
            return None
    except subprocess.TimeoutExpired:
        logger.warning("antiword 执行超时")
        return None
    except Exception as e:
        logger.warning(f"antiword 执行出错: {e}")
        return None


def _extract_doc_text_win32com(file_path: Path) -> str | None:
    """使用 win32com 提取 .doc 文件文本（仅 Windows）"""
    try:
        with com_application("Word.Application") as app:
            doc = app.Documents.Open(str(file_path.resolve()))
            try:
                text = doc.Content.Text
                return text.strip() if text else None
            finally:
                with suppress(Exception):
                    doc.Close()
    except Exception as e:
        logger.warning(f".doc 文本提取失败: {e}", exc_info=True)
        return None


def extract_excel_text(file_path: Path) -> str | None:
    """提取 Excel 表格文本（支持 .xlsx 和 .xls）"""
    extracted_sheets = extract_excel_sheets(file_path)
    if not extracted_sheets:
        return None
    lines: list[str] = []
    for sheet in extracted_sheets:
        if sheet.text:
            lines.append(sheet.text)
    return "\n".join(lines) if lines else None


def extract_excel_sheets(file_path: Path) -> list[ExtractedExcelSheet] | None:
    """提取 Excel 工作簿中每个 Sheet 的文本内容。"""
    suffix = file_path.suffix.lower()

    if suffix == ".xls":
        extracted_sheets = _extract_xls_sheets_xlrd(file_path)
        if extracted_sheets is not None:
            return extracted_sheets
        if _WIN32COM_AVAILABLE:
            return _extract_xls_sheets_win32com(file_path)
        if not _IS_WINDOWS:
            logger.warning("读取 .xls 文件需要 xlrd，请安装: pip install xlrd")
        else:
            logger.warning(
                "读取 .xls 文件失败。请确保 `xlrd` 已安装 (`pip install xlrd`)，或在 Windows 上安装 `pywin32` (`pip install pywin32`) 及 Microsoft Office。"
            )
        return None

    try:
        from openpyxl import load_workbook

        workbook = load_workbook(file_path, read_only=True, data_only=True)
        extracted_sheets: list[ExtractedExcelSheet] = []
        remaining_chars = _EXCEL_PREVIEW_MAX_TOTAL_CHARS
        for worksheet in workbook.worksheets:
            if remaining_chars <= 0:
                break
            extracted_sheet, consumed_chars, stop_after_sheet = _build_excel_sheet_preview(
                sheet_name=worksheet.title,
                rows=worksheet.iter_rows(values_only=True),
                remaining_chars=remaining_chars,
            )
            extracted_sheets.append(extracted_sheet)
            remaining_chars = max(remaining_chars - consumed_chars, 0)
            if stop_after_sheet:
                break
        return extracted_sheets
    except ImportError:
        return None
    except Exception as e:
        logger.warning(f"Excel 文本提取失败: {e}", exc_info=True)
        return None


def _extract_xls_text_xlrd(file_path: Path) -> str | None:
    """使用 xlrd 提取 .xls 文件文本（跨平台）"""
    extracted_sheets = _extract_xls_sheets_xlrd(file_path)
    if not extracted_sheets:
        return None
    lines: list[str] = []
    for sheet in extracted_sheets:
        if sheet.text:
            lines.append(sheet.text)
    return "\n".join(lines) if lines else None


def _extract_xls_text_win32com(file_path: Path) -> str | None:
    """使用 win32com 提取 .xls 文件文本（仅 Windows）"""
    extracted_sheets = _extract_xls_sheets_win32com(file_path)
    if not extracted_sheets:
        return None
    lines: list[str] = []
    for sheet in extracted_sheets:
        if sheet.text:
            lines.append(sheet.text)
    return "\n".join(lines) if lines else None


def _extract_xls_sheets_xlrd(file_path: Path) -> list[ExtractedExcelSheet] | None:
    try:
        import xlrd

        workbook = xlrd.open_workbook(str(file_path))
        extracted_sheets: list[ExtractedExcelSheet] = []
        remaining_chars = _EXCEL_PREVIEW_MAX_TOTAL_CHARS
        for sheet in workbook.sheets():
            if remaining_chars <= 0:
                break
            row_iter = (
                (cell.value for cell in sheet.row(row_idx))
                for row_idx in range(sheet.nrows)
            )
            extracted_sheet, consumed_chars, stop_after_sheet = _build_excel_sheet_preview(
                sheet_name=sheet.name,
                rows=row_iter,
                remaining_chars=remaining_chars,
            )
            extracted_sheets.append(extracted_sheet)
            remaining_chars = max(remaining_chars - consumed_chars, 0)
            if stop_after_sheet:
                break
        return extracted_sheets
    except ImportError:
        return None
    except Exception as e:
        logger.warning(f"xlrd 读取 .xls 失败: {e}")
        return None


def _extract_xls_sheets_win32com(
    file_path: Path,
) -> list[ExtractedExcelSheet] | None:
    try:
        with com_application("Excel.Application") as app:
            workbook = app.Workbooks.Open(str(file_path.resolve()))
            try:
                extracted_sheets: list[ExtractedExcelSheet] = []
                remaining_chars = _EXCEL_PREVIEW_MAX_TOTAL_CHARS
                for worksheet in workbook.Worksheets:
                    if remaining_chars <= 0:
                        break
                    used_range = worksheet.UsedRange
                    row_iter = (
                        (
                            cell.Value
                            for cell in row.Cells
                        )
                        for row in used_range.Rows
                    ) if used_range else ()
                    extracted_sheet, consumed_chars, stop_after_sheet = _build_excel_sheet_preview(
                        sheet_name=str(worksheet.Name),
                        rows=row_iter,
                        remaining_chars=remaining_chars,
                    )
                    extracted_sheets.append(extracted_sheet)
                    remaining_chars = max(remaining_chars - consumed_chars, 0)
                    if stop_after_sheet:
                        break
                return extracted_sheets
            finally:
                with suppress(Exception):
                    workbook.Close(SaveChanges=False)
    except Exception as e:
        logger.warning(f".xls 文本提取失败: {e}", exc_info=True)
        return None


def extract_ppt_text(file_path: Path) -> str | None:
    """提取 PPT 幻灯片文本（支持 .pptx 和 .ppt）"""
    suffix = file_path.suffix.lower()

    # 旧格式 .ppt 需要 win32com
    if suffix == ".ppt":
        return _extract_ppt_text_win32com(file_path)

    # 新格式 .pptx 使用 python-pptx
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except ImportError:
        return None
    except Exception as e:
        logger.warning(f"PPT 文本提取失败: {e}", exc_info=True)
        return None


def _extract_ppt_text_win32com(file_path: Path) -> str | None:
    """使用 win32com 提取 .ppt 文件文本（仅 Windows）"""
    if not _WIN32COM_AVAILABLE:
        if _IS_WINDOWS:
            logger.warning(
                "读取 .ppt 文件需要 pywin32 和 Microsoft PowerPoint，请安装: pip install pywin32"
            )
        else:
            logger.warning(
                "读取 .ppt 旧格式文件仅支持 Windows 环境，请将文件另存为 .pptx 格式"
            )
        return None

    try:
        with com_application("PowerPoint.Application") as app:
            ppt = app.Presentations.Open(str(file_path.resolve()), WithWindow=False)
            try:
                texts = []
                for slide in ppt.Slides:
                    for shape in slide.Shapes:
                        if shape.HasTextFrame:
                            text_frame = shape.TextFrame
                            if text_frame.HasText:
                                texts.append(text_frame.TextRange.Text)
                return "\n".join(texts) if texts else None
            finally:
                with suppress(Exception):
                    ppt.Close()
    except Exception as e:
        logger.warning(f".ppt 文本提取失败: {e}", exc_info=True)
        return None
